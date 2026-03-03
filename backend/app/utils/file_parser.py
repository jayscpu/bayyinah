import fitz  # PyMuPDF


def parse_pdf(file_path: str) -> list[dict]:
    """Extract text from PDF, returning a list of {page_number, text, source_file}."""
    pages = []
    doc = fitz.open(file_path)
    for page_num, page in enumerate(doc, 1):
        text = page.get_text("text")
        if text.strip():
            pages.append({
                "page_number": page_num,
                "text": text.strip(),
                "source_file": file_path.split("/")[-1],
            })
    doc.close()
    return pages


def parse_pptx(file_path: str) -> list[dict]:
    """Extract text from PPTX slides."""
    from pptx import Presentation

    pages = []
    prs = Presentation(file_path)
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
            if hasattr(shape, "notes_slide") and shape.notes_slide:
                for paragraph in shape.notes_slide.notes_text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)

        if texts:
            pages.append({
                "page_number": slide_num,
                "text": "\n".join(texts),
                "source_file": file_path.split("/")[-1],
            })
    return pages
